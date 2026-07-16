"""
src/tools/ppt_export.py
========================
PowerPoint export tool for competitive intelligence briefings.
Generates a professional PPTX presentation from the briefing sections.
"""

from __future__ import annotations

import re
import time
from typing import List, Optional, Tuple, Type

from crewai.tools import BaseTool
from pydantic import BaseModel, Field

from config.settings import settings
from src.utils.logger import get_logger

log = get_logger(__name__)

# Slide colour scheme (McKinsey-inspired)
DARK_BLUE = (22, 33, 62)  # #16213E
MID_BLUE = (15, 52, 96)  # #0F3460
LIGHT_BLUE = (232, 244, 253)  # #E8F4FD
WHITE = (255, 255, 255)
LIGHT_GREY = (248, 249, 250)
TEXT_DARK = (26, 26, 46)


class PPTExportInput(BaseModel):
    markdown_content: str = Field(..., description="Markdown briefing content")
    output_filename: Optional[str] = Field(None, description="Output filename (no extension)")
    run_id: Optional[str] = Field(None, description="Run ID for file naming")
    title: str = Field(
        default="Competitive Intelligence Briefing", description="Presentation title"
    )
    subtitle: str = Field(default="Weekly Strategic Analysis", description="Presentation subtitle")


class PPTExportTool(BaseTool):
    """
    Export the competitive intelligence briefing as a PowerPoint presentation.
    Generates slides for each major section with professional styling.
    """

    name: str = "ppt_export"
    description: str = (
        "Export the briefing as a professional PowerPoint presentation. "
        "Input: full Markdown content. Returns path to generated .pptx file."
    )
    args_schema: Type[BaseModel] = PPTExportInput

    def _run(
        self,
        markdown_content: str,
        output_filename: Optional[str] = None,
        run_id: Optional[str] = None,
        title: str = "Competitive Intelligence Briefing",
        subtitle: str = "Weekly Strategic Analysis",
    ) -> str:
        """Generate PPTX from Markdown briefing."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return "python-pptx not installed. Run: pip install python-pptx"

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # ── Cover slide ──────────────────────────────────────
        self._add_cover_slide(prs, title, subtitle)

        # ── Parse Markdown sections ───────────────────────────
        sections = self._parse_sections(markdown_content)
        for section_title, section_content in sections:
            self._add_content_slide(prs, section_title, section_content)

        # ── Thank you / End slide ─────────────────────────────
        self._add_end_slide(prs)

        # Save
        filename = output_filename or (
            f"briefing_{run_id}" if run_id else f"briefing_{int(time.time())}"
        )
        output_path = settings.output_dir / f"{filename}.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        size_kb = output_path.stat().st_size // 1024
        log.info(f"[PPT] Generated: {output_path} ({size_kb} KB)")
        return f"PowerPoint exported: {output_path} ({size_kb} KB)"

    def _rgb(self, tup: Tuple[int, int, int]):
        from pptx.dml.color import RGBColor

        return RGBColor(*tup)

    def _add_cover_slide(self, prs, title: str, subtitle: str) -> None:
        from pptx.util import Inches, Pt

        blank_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(blank_layout)

        # Background
        self._fill_background(slide, DARK_BLUE)

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self._rgb(WHITE)

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = self._rgb(LIGHT_BLUE)

        # Accent line
        from pptx.util import Inches

        line = slide.shapes.add_shape(1, Inches(1.5), Inches(4.0), Inches(10), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(MID_BLUE)
        line.line.fill.background()

    def _add_content_slide(self, prs, section_title: str, content: str) -> None:
        from pptx.util import Inches, Pt

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Light background
        self._fill_background(slide, WHITE)

        # Title bar
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.9))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self._rgb(DARK_BLUE)

        # Accent line under title
        line = slide.shapes.add_shape(1, Inches(0.3), Inches(1.15), Inches(12.73), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(MID_BLUE)
        line.line.fill.background()

        # Content
        clean_content = self._clean_markdown(content)
        # Limit content to fit slide
        if len(clean_content) > 1500:
            clean_content = clean_content[:1497] + "..."

        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(12.73), Inches(5.8))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = clean_content
        p2.font.size = Pt(11)
        p2.font.color.rgb = self._rgb(TEXT_DARK)

    def _add_end_slide(self, prs) -> None:
        from pptx.util import Inches, Pt

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        self._fill_background(slide, DARK_BLUE)
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "End of Competitive Intelligence Briefing"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self._rgb(WHITE)

    def _fill_background(self, slide, color: Tuple[int, int, int]) -> None:
        """Fill slide background with a solid colour."""

        background = slide.shapes.add_shape(
            1,
            0,
            0,
            slide.shapes._spTree.getparent().getparent().cxSp,
            slide.shapes._spTree.getparent().getparent().cySp,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self._rgb(color)
        background.line.fill.background()

    def _parse_sections(self, markdown: str) -> List[Tuple[str, str]]:
        """Parse Markdown H2 sections."""
        sections = []
        current_title = "Overview"
        current_content = []

        for line in markdown.split("\n"):
            if line.startswith("## "):
                if current_content:
                    sections.append((current_title, "\n".join(current_content)))
                current_title = line[3:].strip()
                current_content = []
            elif not line.startswith("# "):
                current_content.append(line)

        if current_content:
            sections.append((current_title, "\n".join(current_content)))

        # Limit to 15 slides (cover + sections + end)
        return sections[:13]

    def _clean_markdown(self, text: str) -> str:
        """Remove Markdown syntax for plain text slides."""
        text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
        text = re.sub(r"\*(.*?)\*", r"\1", text)
        text = re.sub(r"#{1,6}\s", "", text)
        text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
        text = re.sub(r"\[(\d+)\]", r"[\1]", text)
        text = re.sub(r"^\s*[-*+]\s", "• ", text, flags=re.MULTILINE)
        text = re.sub(r"\|", " | ", text)
        text = re.sub(r"---+", "", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


# Singleton
ppt_export_tool = PPTExportTool()

__all__ = ["PPTExportTool", "ppt_export_tool"]
